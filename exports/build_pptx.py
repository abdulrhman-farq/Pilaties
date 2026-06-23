# -*- coding: utf-8 -*-
"""
توليد عرض PowerPoint للمستثمر — استوديو Reformer Pilates الرياض.
يضبط اتجاه الفقرات RTL ومحاذاتها لليمين. شغّل: python3 build_pptx.py
ملاحظة: PowerPoint/Keynote يشكّلان الحروف العربية تلقائياً عند الفتح.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1F, 0x4E, 0x78)
TEAL = RGBColor(0x2E, 0x8B, 0x8B)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
LGREY = RGBColor(0xF0, 0xF3, 0xF8)
GREEN = RGBColor(0xC6, 0xEF, 0xCE)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def set_font(run, size=None, bold=None, color=None):
    """يضبط خط النص اللاتيني + العربي (complex-script) لضمان وصل الحروف العربية."""
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = FONT  # a:latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):  # خط شرق آسيوي + خط معقّد (العربي)
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=18, color=DARK, bold=False,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        rtl(p)
        r = p.add_run(); r.text = line
        set_font(r, size=Pt(size), bold=bold, color=color)
    return tb


def title_bar(s, t, idx=None):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Pt(4), TEAL)
    text(s, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.8), t,
         size=26, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if idx is not None:
        text(s, Inches(0.3), Inches(0.3), Inches(1.2), Inches(0.6), str(idx),
             size=18, color=GOLD, bold=True, align=PP_ALIGN.LEFT)


def bullets(s, items, x=Inches(0.7), y=Inches(1.6), w=Inches(12), size=20,
            color=DARK, gap=True):
    tb = s.shapes.add_textbox(x, y, w, SH - y - Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT; p.line_spacing = 1.2
        if gap:
            p.space_after = Pt(10)
        rtl(p)
        r = p.add_run(); r.text = "•  " + it
        set_font(r, size=Pt(size), color=color)
    return tb


def table(s, rows, x, y, w, h, header=True, col_widths=None,
          highlight=None, font_size=16):
    highlight = highlight or {}
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, x, y, w, h).table
    if col_widths:
        for c, cw in enumerate(col_widths):
            gt.columns[c].width = cw
    for r in range(nr):
        for c in range(nc):
            cell = gt.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            if r == 0 and header:
                cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
            elif r in highlight:
                cell.fill.solid(); cell.fill.fore_color.rgb = highlight[r]
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LGREY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            rtl(p)
            run = p.add_run(); run.text = str(rows[r][c])
            set_font(run, size=Pt(font_size),
                     bold=(r == 0) or (c == 0 and nc > 2),
                     color=WHITE if r == 0 else DARK)
    return gt


# ============================================================
# الشريحة 0 — الغلاف
# ============================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.0), SW, Inches(1.5), TEAL)
text(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(1.2),
     "استوديو Reformer Pilates نسائي — الرياض", size=34, color=WHITE,
     bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.2),
     "عرض استثماري", size=30, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.0),
     "120 م²  |  6 أجهزة Reformer  |  دراسة جدوى — يونيو 2026",
     size=20, color=WHITE, align=PP_ALIGN.CENTER)
text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
     "وثيقة مُعدّة لعرضها على مستثمر / صندوق استثماري",
     size=13, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

# 1 — الفرصة
s = slide(); title_bar(s, "الفرصة", 1)
bullets(s, [
    "استوديو بيلاتس نسائي بوتيكي في الرياض — في صميم أسرع شريحتين نمواً بالسوق: النساء + القطاع البوتيكي.",
    "سوق اللياقة السعودي: ~1.5 مليار دولار (2025) ← ~3 مليار دولار (2030).",
    "نمو القطاع البوتيكي ~13.6% سنوياً | نمو شريحة النساء ~13.25% سنوياً.",
    "مدعوم برؤية 2030 (مستهدف 40% ممارسة للنشاط البدني).",
], size=21)

# 2 — المشكلة/الفجوة
s = slide(); title_bar(s, "المشكلة والفجوة", 2)
bullets(s, [
    "فاخر جداً: حصة بـ 300–375 ريال — حاجز سعري مرتفع.",
    "اقتصادي مزدحم: أندية ضخمة تفقد الطابع الحميمي والخصوصية.",
    "الفجوة: استوديو بوتيكي راقٍ، نسائي بالكامل، بسعر في المتناول (150 ريال) وتجربة شبه شخصية (6 عميلات كحد أقصى).",
], size=22)

# 3 — الحل
s = slide(); title_bar(s, "الحل", 3)
bullets(s, [
    "استوديو 120 م² بـ 6 أجهزة Reformer (قابلة للتوسعة إلى 8).",
    "تجربة حميمية فاخرة وخصوصية نسائية تامة.",
    "موقع راقٍ + هوية علامة قوية.",
    "نموذج إيرادات متعدد: حصص مفردة + باقات + اشتراكات + حصص خاصة.",
], size=22)

# 4 — السوق المستهدف
s = slide(); title_bar(s, "السوق المستهدف", 4)
bullets(s, [
    "سيدات 25–45 سنة، دخل متوسط إلى مرتفع.",
    "العاملات (الصباح الباكر والمساء) + ربات المنزل (فترات النهار).",
    "محرّكات الطلب: الوعي الصحي، نمط الحياة، ما بعد الولادة، التعافي.",
], size=22)

# 5 — نموذج العمل
s = slide(); title_bar(s, "نموذج العمل", 5)
table(s, [
    ["المصدر", "السعر"],
    ["حصة مفردة", "150 ريال"],
    ["باقة 8 حصص", "1,000 ريال (125/حصة)"],
    ["اشتراك شهري غير محدود", "1,600 ريال"],
    ["حصة خاصة (Private/Duet)", "250–400 ريال"],
], Inches(3.2), Inches(1.6), Inches(7), Inches(3),
    col_widths=[Inches(3.5), Inches(3.5)], font_size=18)
text(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.8),
     "الطاقة: 6 أجهزة × 7 حصص × 26 يوم = 1,092 حصة-عميلة/شهر",
     size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# 6 — الاقتصاديات
s = slide(); title_bar(s, "الاقتصاديات", 6)
table(s, [
    ["المؤشر", "القيمة"],
    ["إجمالي الاستثمار", "275,000 ريال"],
    ["تكلفة تشغيل شهرية", "61,000 ريال"],
    ["نقطة التعادل", "37% إشغال فقط (~2.2 عميلة/حصة)"],
], Inches(2.7), Inches(1.7), Inches(8), Inches(2.6),
    col_widths=[Inches(4), Inches(4)], font_size=18)
text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.8),
     "رأسمال منخفض + نقطة تعادل منخفضة = مخاطرة تشغيلية محدودة",
     size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

# 7 — التوقعات المالية
s = slide(); title_bar(s, "التوقعات المالية (سيناريوهات الإشغال)", 7)
table(s, [
    ["السيناريو", "الإيراد الشهري", "صافي شهري", "صافي سنوي"],
    ["إشغال 30%", "49,140", "(11,860)", "(142,320)"],
    ["إشغال 50%", "81,900", "20,900", "250,800"],
    ["إشغال 70%", "114,660", "53,660", "643,920"],
    ["إشغال 90%", "147,420", "86,420", "1,037,040"],
], Inches(1.2), Inches(1.7), Inches(11), Inches(3.8),
    highlight={2: GREEN}, font_size=17)

# 8 — مؤشرات العائد
s = slide(); title_bar(s, "مؤشرات العائد", 8)
table(s, [
    ["المؤشر", "القيمة"],
    ["NPV (5 سنوات @15%)", "~ +653,000 ريال"],
    ["IRR (معدل العائد الداخلي)", "~ 66%"],
    ["ROI (إشغال 50%)", "~ 91% سنوياً"],
    ["فترة الاسترداد", "~ 13 شهر"],
    ["أرباح أول 3 سنوات", "~ 666,000 ريال"],
], Inches(3.0), Inches(1.6), Inches(7.3), Inches(3.4),
    col_widths=[Inches(3.8), Inches(3.5)], font_size=18)

# 9 — خطة النمو
s = slide(); title_bar(s, "خطة النمو", 9)
bullets(s, [
    "السنة 1: بناء قاعدة العميلات (الوصول لإشغال 50%+).",
    "السنة 2: تثبيت الربحية + إضافة جهازين (120 م² تستوعب 8 أجهزة بنفس الإيجار).",
    "السنة 3+: غرفة نشاط ثانٍ (Mat/Barre) أو فرع ثانٍ عند إشغال 80%+ مستقر.",
], size=22)

# 10 — المخاطر
s = slide(); title_bar(s, "المخاطر وتخفيفها", 10)
table(s, [
    ["المخاطرة", "إجراء التخفيف"],
    ["المنافسة", "تموضع بوتيكي + ولاء للعلامة لا للسعر"],
    ["انخفاض الإشغال", "نقطة تعادل منخفضة (37%) + إحالات + باقات"],
    ["الاعتماد على مدربة", "تعدد المدربات + توثيق المنهجية + بناء علامة"],
    ["ارتفاع الإيجار", "عقد طويل بسقف زيادة متفق عليه"],
], Inches(1.2), Inches(1.7), Inches(11), Inches(3.6),
    col_widths=[Inches(3.5), Inches(7.5)], font_size=17)

# 11 — الطلب
s = slide(); title_bar(s, "الطلب التمويلي", 11)
text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.0),
     "التمويل المطلوب: 300,000 ريال (يشمل هامش أمان ~25 ألف)",
     size=26, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
bullets(s, [
    "استرداد كامل خلال السنة الثانية.",
    "نقطة مراجعة أداء عند الشهر السادس قبل أي قرار توسعة.",
    "التوصية: المضي قدماً (GO).",
], y=Inches(3.0), size=22)
rect(s, Inches(2.5), Inches(5.6), Inches(8.3), Inches(1.1), TEAL)
text(s, Inches(2.5), Inches(5.6), Inches(8.3), Inches(1.1),
     "استثمار صغير، نقطة تعادل منخفضة، عائد استثنائي — في سوق صاعد",
     size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
     anchor=MSO_ANCHOR.MIDDLE)

# 12 — لماذا الآن
s = slide(); title_bar(s, "لماذا الآن؟", 12)
bullets(s, [
    "السوق في ذروة النمو ومدعوم حكومياً (رؤية 2030).",
    "شريحة النساء هي الأسرع نمواً وغير مُشبَعة في كثير من الأحياء.",
    "رأسمال صغير + عائد داخلي يفوق 65% = نسبة عائد/مخاطرة جذابة جداً.",
], size=22)
text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.0),
     "« استثمار صغير، نقطة تعادل منخفضة، عائد استثنائي، في سوق صاعد »",
     size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "عرض-المستثمر-بيلاتس.pptx")
prs.save(out)
print("تم إنشاء:", out, "—", len(prs.slides._sldIdLst), "شريحة")
